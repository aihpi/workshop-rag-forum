"""Build the 'Introducing turbovec' deck (.pptx) on top of the AISC template.

Loads the AI Service Centre template presentation (so the deck inherits its master,
theme, fonts and footer/page-number chrome), drops the template's own slides, and adds
eight content slides using the template's 'Inhalt 1 - <colour>' layouts:

  1. The problem
  2-5. How turbovec works (four-step, layman-friendly walk-through)
  6. Wikipedia results - size      (theme-coloured bar chart)
  7. Wikipedia results - latency   (theme-coloured bar chart)
  8. Three open questions

All numbers come from data/results.json (03_compress_turbovec.py + 05_benchmark_latency.py).
Charts are rendered with matplotlib and embedded as transparent PNGs. No PDF is produced
(create one from the .pptx with PowerPoint / LibreOffice if needed).

Template path: env TEMPLATE_PPTX, else tmp/20260529_on-premis_ai_for_public_knowledge.pptx
"""

from __future__ import annotations

import io
import json
import os
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
HERE = Path(__file__).resolve().parent
RESULTS = REPO_ROOT / "data" / "results.json"
TEMPLATE = Path(os.environ.get(
    "TEMPLATE_PPTX",
    REPO_ROOT / "tmp" / "20260529_on-premis_ai_for_public_knowledge.pptx"))
PPTX_OUT = HERE / "turbovec_intro.pptx"

# AISC theme palette (from the template's theme1.xml).
GREY = "#5A6065"
ORANGE = "#DD6108"
DARK = "#2B2F33"

# Layout indices in the template (see 'Inhalt 1 - <colour>').
L_RED, L_ORANGE, L_BLUE, L_GREEN = 3, 2, 4, 5

R = json.loads(RESULTS.read_text())
TV = R["turbovec"]
LAT = R["latency"]
F32_MB = R["baseline_float32"]["bytes"] / 1e6
BITS = sorted((int(k[0]) for k in TV if k.endswith("bit")), reverse=True)  # 4,3,2


# --------------------------------------------------------------------- charts
def _bar_png(values, labels, fmt, title, ylabel) -> io.BytesIO:
    fig = plt.figure(figsize=(5.3, 4.3), dpi=200)
    fig.patch.set_alpha(0.0)
    ax = fig.add_axes([0.16, 0.12, 0.80, 0.74])
    ax.patch.set_alpha(0.0)
    colors = [GREY] + [ORANGE] * (len(values) - 1)
    bars = ax.bar(labels, values, color=colors, width=0.62, zorder=3)
    ax.set_title(title, fontsize=15, weight="bold", color=DARK, pad=12)
    ax.set_ylabel(ylabel, fontsize=12, color=GREY)
    ax.grid(axis="y", color="#DDDFE0", zorder=0)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    for s in ("left", "bottom"):
        ax.spines[s].set_color("#DDDFE0")
    ax.tick_params(labelsize=12, colors=DARK)
    ax.set_ylim(0, max(values) * 1.18)
    for rect, v in zip(bars, values):
        ax.text(rect.get_x() + rect.get_width() / 2, v, fmt(v), ha="center",
                va="bottom", fontsize=12, weight="bold", color=DARK)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def size_chart():
    vals = [F32_MB] + [TV[f"{b}bit"]["bytes"] / 1e6 for b in BITS]
    labs = ["float32"] + [f"{b}-bit" for b in BITS]
    return _bar_png(vals, labs, lambda v: f"{v:.0f} MB", "Index size on disk",
                    "megabytes")


def latency_chart():
    vals = [LAT["float32_ms_per_query"]] + [TV[f"{b}bit"]["latency_ms_per_query"]
                                            for b in BITS]
    labs = ["float32"] + [f"{b}-bit" for b in BITS]
    return _bar_png(vals, labs, lambda v: f"{v:.2f}", "Search latency (ms / query)",
                    "ms per query")


# ------------------------------------------------------------------ pptx build
def clear_slides(prs: Presentation) -> None:
    """Remove every slide the template ships with, leaving masters/layouts intact."""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def _set_paragraph(p, text: str, size=None):
    # Bold the leading label up to the first ': ' or ' - ' (em dash) separator.
    p.clear()
    lead, rest = text, ""
    for sep in (": ", " — "):
        i = text.find(sep)
        if i != -1:
            lead, rest = text[: i + len(sep)], text[i + len(sep):]
            break
    r1 = p.add_run(); r1.text = lead; r1.font.bold = True
    if rest:
        r2 = p.add_run(); r2.text = rest; r2.font.bold = False
    if size:
        for run in p.runs:
            run.font.size = Pt(size)


def fill_body(ph, bullets, size=None):
    tf = ph.text_frame
    tf.word_wrap = True
    for i, btxt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        _set_paragraph(p, btxt, size=size)


def add_slide(prs, layout_idx, title, bullets, chart=None, body_size=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    if chart is not None:
        # Narrow the body to the left half; chart fills the right half.
        body.left, body.top = Inches(0.5), Inches(1.7)
        body.width, body.height = Inches(6.2), Inches(5.0)
        fill_body(body, bullets, size=body_size or 16)
        slide.shapes.add_picture(chart, Inches(7.05), Inches(1.95),
                                 width=Inches(5.7))
    else:
        fill_body(body, bullets, size=body_size)
    return slide


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    add_slide(prs, L_RED,
              "The problem: embeddings are big — and big is slow", [
                  "Each text chunk is stored as one float32 vector: 384 numbers "
                  "× 4 bytes = 1.5 KB.",
                  "Our 100k-chunk slice of German Wikipedia already needs "
                  f"{F32_MB:.0f} MB; 100M chunks would be ~150 GB.",
                  "These vectors must sit in RAM for fast search — so memory "
                  "is both the bottleneck and the bill.",
                  f"Exact float32 search already costs "
                  f"{LAT['float32_ms_per_query']:.1f} ms/query "
                  f"(~{LAT['float32_qps']:.0f} q/s) on just 100k vectors.",
                  "Can we shrink the vectors without wrecking retrieval quality? "
                  "— that is what turbovec does.",
              ])

    add_slide(prs, L_BLUE,
              "How it works (1/4): an embedding is just a list of numbers", [
                  "An embedding turns a piece of text into ~384 numbers — "
                  "coordinates that place similar meanings close together.",
                  "Search means: find the chunks whose coordinates sit nearest to "
                  "the question's coordinates.",
                  "Today each number is a float32: 32 bits, about 7 digits of "
                  "precision.",
                  "That is like writing every measurement to 7 decimals when you "
                  "only need one — most of those bits never change the answer.",
              ])

    add_slide(prs, L_BLUE,
              "How it works (2/4): round the numbers into a few buckets", [
                  "turbovec keeps only a handful of levels per number: 2 bits = 4 "
                  "levels, 3 bits = 8, 4 bits = 16.",
                  "Each number is snapped to its nearest level — like naming a "
                  "colour from a 16-crayon box instead of millions of shades.",
                  "A per-vector scale places those few levels where that vector's "
                  "values actually lie, so little is lost.",
                  "32 bits → 2 bits is 16× fewer bits per number; the whole "
                  "index shrinks almost as much. This is called quantization.",
              ])

    add_slide(prs, L_BLUE,
              "How it works (3/4): why rounding doesn't break search", [
                  "To rank nearest neighbours you only need relative positions, not "
                  "exact values.",
                  "Rounding nudges every point a little but keeps it in roughly the "
                  "same place — neighbours stay neighbours.",
                  "So the right answers still land near the top: at 4-bit we keep "
                  f"{TV['4bit']['recall@10']*100:.0f}% of them, at 2-bit "
                  f"{TV['2bit']['recall@10']*100:.0f}%.",
                  "Want the last bit of quality back? Re-rank just the top hits "
                  "using the original float32 vectors.",
              ])

    add_slide(prs, L_BLUE,
              "How it works (4/4): why smaller is also faster", [
                  "Searching means reading every stored vector and comparing it to "
                  "the query — work that is limited mostly by memory speed.",
                  "Smaller codes mean far less data to move, so search speeds up "
                  "roughly in step with the shrink.",
                  "turbovec compares the packed codes directly — it never "
                  "unpacks them back to float32 first.",
                  "It leans on special CPU instructions (SIMD: AVX-512 on x86, NEON "
                  "on ARM) that crunch many numbers at once.",
              ])

    size_rows = [f"100k chunks, dim 384.  Baseline float32 = {F32_MB:.0f} MB."]
    for b in BITS:
        e = TV[f"{b}bit"]
        size_rows.append(f"{b}-bit: {e['bytes']/1e6:.1f} MB "
                         f"({e['compression_ratio']:.1f}× smaller), "
                         f"recall@10 = {e['recall@10']:.3f}")
    size_rows.append("More bits → better recall; fewer bits → smaller index.")
    add_slide(prs, L_GREEN, "Results on Wikipedia (1/2): up to 15× smaller",
              size_rows, chart=size_chart(), body_size=15)

    lat_rows = [f"Exact float32 search: {LAT['float32_ms_per_query']:.2f} ms/query "
                f"(~{LAT['float32_qps']:.0f} q/s)."]
    for b in BITS:
        e = TV[f"{b}bit"]
        lat_rows.append(f"{b}-bit: {e['latency_ms_per_query']:.2f} ms/query "
                        f"({e['speedup_vs_float32']:.1f}× faster, "
                        f"~{e['qps']:.0f} q/s)")
    lat_rows.append("Speed-up tracks compression: search is memory-bandwidth bound, "
                    "so smaller codes run faster.")
    add_slide(prs, L_GREEN, "Results on Wikipedia (2/2): up to 16× faster",
              lat_rows, chart=latency_chart(), body_size=15)

    add_slide(prs, L_ORANGE, "Three questions this raises", [
        f"How low can we go?  Is 2-bit's ~{TV['2bit']['recall@10']:.2f} recall good "
        "enough, or do we re-rank the top candidates with float32?",
        "Does it hold on our data and at scale?  These are 100k German-Wikipedia "
        "chunks with minilm — what about 10M+ vectors and our domain "
        "embeddings (qwen3 / octen)?",
        "Where does it fit in the stack?  Drop-in for our vector DB or a separate "
        "index? How does it compose with query transformation, MCP tools and "
        "ingestion?",
    ])

    prs.save(str(PPTX_OUT))
    print(f"Wrote {PPTX_OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
